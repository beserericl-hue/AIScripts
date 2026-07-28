"""Office-format (xlsx / pptx / docx) → text extractors for the bulk
supporting-evidence flow.

The bulk-evidence uploader lets a coordinator drop a folder of supporting
files (agendas, syllabi, enrollment spreadsheets, program decks) that the
self-study REFERENCES but that carry no standard/sub-standard number. To
recommend a placement we need the file's plain text so it can be (a) matched
against the imported narratives by name/content and (b) fed to the
``/ai/placement/recommend`` SpecMatcher.

pypdf (PDF) and mammoth (DOCX) already live in the service; this module adds
the two missing Office formats with light-weight, pure-Python readers:

* xlsx  → ``openpyxl`` (read-only, values only — no formulas / styles)
* pptx  → ``python-pptx`` (every shape's text frame + table cells + notes)

Both raise ``ValueError`` on malformed / oversized input so the route handler
can surface a 400 rather than a 502.
"""
from __future__ import annotations

import io


_MAX_BYTES = 50 * 1024 * 1024  # 50 MB — matches the upload cap.
# Cap the extracted text so a giant workbook can't blow the placement prompt.
_MAX_TEXT_CHARS = 200_000


def _guard(data: bytes, label: str) -> None:
    if not data:
        raise ValueError(f"{label} is empty")
    if len(data) > _MAX_BYTES:
        raise ValueError(
            f"{label} exceeds {_MAX_BYTES // (1024 * 1024)} MB cap "
            f"(got {len(data) // (1024 * 1024)} MB)"
        )


def extract_text_from_xlsx_bytes(data: bytes) -> str:
    """Flatten every sheet of an xlsx to tab-separated rows.

    Sheet names become headers so a placement match can see e.g. an
    "Enrollment" or "Graduation Rates" tab title, which is often the strongest
    signal for which standard a workbook supports.
    """
    _guard(data, "xlsx")
    from openpyxl import load_workbook

    try:
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001 — openpyxl raises many types
        raise ValueError(f"could not read xlsx: {exc}") from exc

    out: list[str] = []
    try:
        for ws in wb.worksheets:
            out.append(f"# Sheet: {ws.title}")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    out.append("\t".join(cells))
                if sum(len(s) for s in out) > _MAX_TEXT_CHARS:
                    break
    finally:
        wb.close()

    return "\n".join(out)[:_MAX_TEXT_CHARS]


def extract_text_from_pptx_bytes(data: bytes) -> str:
    """Pull every slide's text (shapes, tables, speaker notes) as markdown.

    Slides are separated by a blank line and headed ``## Slide N`` so the
    downstream chunker/matcher can use slide boundaries as soft cuts and see
    slide titles.
    """
    _guard(data, "pptx")
    from pptx import Presentation

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001 — python-pptx raises many types
        raise ValueError(f"could not read pptx: {exc}") from exc

    out: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        out.append(f"## Slide {i}")
        for shape in slide.shapes:
            try:
                if shape.has_table:
                    for r in shape.table.rows:
                        cells = [c.text.strip() for c in r.cells if c.text.strip()]
                        if cells:
                            out.append("\t".join(cells))
                    continue
            except Exception:  # noqa: BLE001 — shape.has_table can raise on odd shapes
                pass
            if getattr(shape, "has_text_frame", False):
                txt = shape.text_frame.text.strip()
                if txt:
                    out.append(txt)
        # Speaker notes frequently spell out the narrative context.
        try:
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    out.append(f"(notes) {notes}")
        except Exception:  # noqa: BLE001
            pass
        if sum(len(s) for s in out) > _MAX_TEXT_CHARS:
            break

    return "\n\n".join(out)[:_MAX_TEXT_CHARS]


def extract_text_from_docx_bytes(data: bytes) -> str:
    """DOCX → plain text via mammoth (already used by the import path)."""
    _guard(data, "docx")
    import mammoth

    try:
        result = mammoth.extract_raw_text(io.BytesIO(data))
    except Exception as exc:  # noqa: BLE001
        raise ValueError(f"could not read docx: {exc}") from exc
    return (result.value or "")[:_MAX_TEXT_CHARS]


# mimeType / extension → extractor. The route handler resolves by whichever it
# has; extension wins when the browser sends a generic octet-stream.
def extract_text(data: bytes, *, mime_type: str = "", filename: str = "") -> str:
    """Dispatch to the right extractor by mimeType then filename extension.

    Returns the extracted text (possibly empty). Raises ``ValueError`` for an
    unsupported type so the caller can 415 rather than silently returning "".
    """
    mt = (mime_type or "").lower()
    fn = (filename or "").lower()

    def is_(kind: str, *exts: str) -> bool:
        return kind in mt or any(fn.endswith(e) for e in exts)

    if is_("spreadsheetml", ".xlsx"):
        return extract_text_from_xlsx_bytes(data)
    if is_("presentationml", ".pptx"):
        return extract_text_from_pptx_bytes(data)
    if is_("wordprocessingml", ".docx"):
        return extract_text_from_docx_bytes(data)
    if is_("pdf", ".pdf"):
        from .pdf_extract import extract_text_from_pdf_bytes

        return extract_text_from_pdf_bytes(data)
    raise ValueError(f"unsupported file type for text extraction: mime={mime_type!r} name={filename!r}")
